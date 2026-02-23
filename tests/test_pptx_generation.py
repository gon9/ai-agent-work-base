"""
PptxGenerationSkillのテスト。

正常系・異常系を網羅する。
"""

from __future__ import annotations

import json
import os
import tempfile

import pytest
from pptx import Presentation

from ai_agent_work_base.skills.pptx_generation import PptxGenerationSkill


@pytest.fixture
def skill() -> PptxGenerationSkill:
    """テスト用スキルインスタンス。"""
    return PptxGenerationSkill()


@pytest.fixture
def tmp_pptx(tmp_path) -> str:
    """一時出力ファイルパス。"""
    return str(tmp_path / "test.pptx")


# ── スライド定義サンプル ────────────────────────────────────────────────────────

TITLE_SLIDE = {"type": "title", "title": "テストタイトル", "subtitle": "サブタイトル", "notes": "発表者ノート"}

BULLETS_SLIDE = {
    "type": "bullets",
    "title": "箇条書きスライド",
    "bullets": ["項目A", "項目B", {"level": 1, "text": "サブ項目"}],
    "notes": "ノート",
}

CONTENT_SLIDE = {"type": "content", "title": "コンテンツスライド", "body": "本文テキストです。"}

CHART_SLIDE = {
    "type": "chart",
    "title": "棒グラフ",
    "chart": {
        "chart_type": "bar",
        "categories": ["Q1", "Q2", "Q3", "Q4"],
        "series": [{"name": "売上", "values": [100, 120, 90, 150]}],
    },
}

TABLE_SLIDE = {
    "type": "table",
    "title": "テーブル",
    "table": {
        "headers": ["項目", "値", "備考"],
        "rows": [["A", "100", "良好"], ["B", "80", "普通"]],
    },
}

TWO_COLUMN_SLIDE = {
    "type": "two_column",
    "title": "2カラム比較",
    "columns": [
        {"title": "メリット", "bullets": ["速い", "安い"]},
        {"title": "デメリット", "bullets": ["難しい", "複雑"]},
    ],
}

ALL_SLIDES = [TITLE_SLIDE, BULLETS_SLIDE, CONTENT_SLIDE, CHART_SLIDE, TABLE_SLIDE, TWO_COLUMN_SLIDE]


# ── 正常系テスト ──────────────────────────────────────────────────────────────

class TestPptxGenerationSkillNormal:
    """正常系テスト。"""

    def test_skill_name(self, skill: PptxGenerationSkill) -> None:
        """スキル名が正しいこと。"""
        assert skill.name == "generate_pptx"

    def test_skill_description_not_empty(self, skill: PptxGenerationSkill) -> None:
        """説明が空でないこと。"""
        assert len(skill.description) > 0

    def test_parameters_schema(self, skill: PptxGenerationSkill) -> None:
        """パラメータスキーマにrequiredフィールドが含まれること。"""
        params = skill.parameters
        assert params["type"] == "object"
        assert "title" in params["properties"]
        assert "slides" in params["properties"]
        assert "file_path" in params["properties"]
        assert set(params["required"]) == {"title", "slides", "file_path"}

    def test_generate_title_slide(self, skill: PptxGenerationSkill, tmp_pptx: str) -> None:
        """タイトルスライド1枚を生成できること。"""
        result = skill.execute(
            title="テスト",
            slides=json.dumps([TITLE_SLIDE]),
            file_path=tmp_pptx,
        )
        assert "生成しました" in result
        assert os.path.exists(tmp_pptx)
        prs = Presentation(tmp_pptx)
        assert len(prs.slides) == 1

    def test_generate_all_slide_types(self, skill: PptxGenerationSkill, tmp_pptx: str) -> None:
        """全スライドタイプを含む複数スライドを生成できること。"""
        result = skill.execute(
            title="全タイプテスト",
            slides=json.dumps(ALL_SLIDES),
            file_path=tmp_pptx,
        )
        assert "生成しました" in result
        assert os.path.exists(tmp_pptx)
        prs = Presentation(tmp_pptx)
        assert len(prs.slides) == len(ALL_SLIDES)

    def test_generate_with_footer(self, skill: PptxGenerationSkill, tmp_pptx: str) -> None:
        """フッターを指定して生成できること。"""
        result = skill.execute(
            title="フッターテスト",
            slides=json.dumps([TITLE_SLIDE, CONTENT_SLIDE]),
            file_path=tmp_pptx,
            footer="テスト | AI Agent",
        )
        assert "生成しました" in result
        assert os.path.exists(tmp_pptx)

    def test_slides_as_list(self, skill: PptxGenerationSkill, tmp_pptx: str) -> None:
        """slidesをリストで渡しても生成できること。"""
        result = skill.execute(
            title="リスト渡しテスト",
            slides=[TITLE_SLIDE, BULLETS_SLIDE],
            file_path=tmp_pptx,
        )
        assert "生成しました" in result
        assert os.path.exists(tmp_pptx)

    def test_slides_with_code_block(self, skill: PptxGenerationSkill, tmp_pptx: str) -> None:
        """LLMがコードブロックで返したJSON文字列でも生成できること。"""
        slides_str = f"```json\n{json.dumps([TITLE_SLIDE])}\n```"
        result = skill.execute(
            title="コードブロックテスト",
            slides=slides_str,
            file_path=tmp_pptx,
        )
        assert "生成しました" in result
        assert os.path.exists(tmp_pptx)

    def test_output_directory_created(self, skill: PptxGenerationSkill, tmp_path: str) -> None:
        """出力ディレクトリが存在しなくても自動作成されること。"""
        nested_path = str(tmp_path / "nested" / "dir" / "output.pptx")
        result = skill.execute(
            title="ディレクトリ作成テスト",
            slides=json.dumps([TITLE_SLIDE]),
            file_path=nested_path,
        )
        assert "生成しました" in result
        assert os.path.exists(nested_path)

    def test_chart_slide_bar(self, skill: PptxGenerationSkill, tmp_pptx: str) -> None:
        """棒グラフスライドが生成できること。"""
        result = skill.execute(
            title="棒グラフテスト",
            slides=json.dumps([TITLE_SLIDE, CHART_SLIDE]),
            file_path=tmp_pptx,
        )
        assert "生成しました" in result
        prs = Presentation(tmp_pptx)
        assert len(prs.slides) == 2

    def test_chart_slide_line(self, skill: PptxGenerationSkill, tmp_pptx: str) -> None:
        """折れ線グラフスライドが生成できること。"""
        line_chart = {
            "type": "chart",
            "title": "折れ線グラフ",
            "chart": {
                "chart_type": "line",
                "categories": ["1月", "2月", "3月"],
                "series": [
                    {"name": "A", "values": [10, 20, 15]},
                    {"name": "B", "values": [5, 12, 18]},
                ],
            },
        }
        result = skill.execute(
            title="折れ線テスト",
            slides=json.dumps([line_chart]),
            file_path=tmp_pptx,
        )
        assert "生成しました" in result

    def test_chart_slide_pie(self, skill: PptxGenerationSkill, tmp_pptx: str) -> None:
        """円グラフスライドが生成できること。"""
        pie_chart = {
            "type": "chart",
            "title": "円グラフ",
            "chart": {
                "chart_type": "pie",
                "categories": ["A", "B", "C"],
                "series": [{"name": "シェア", "values": [40, 35, 25]}],
            },
        }
        result = skill.execute(
            title="円グラフテスト",
            slides=json.dumps([pie_chart]),
            file_path=tmp_pptx,
        )
        assert "生成しました" in result

    def test_table_slide(self, skill: PptxGenerationSkill, tmp_pptx: str) -> None:
        """テーブルスライドが生成できること。"""
        result = skill.execute(
            title="テーブルテスト",
            slides=json.dumps([TABLE_SLIDE]),
            file_path=tmp_pptx,
        )
        assert "生成しました" in result
        prs = Presentation(tmp_pptx)
        assert len(prs.slides) == 1

    def test_two_column_slide(self, skill: PptxGenerationSkill, tmp_pptx: str) -> None:
        """2カラムスライドが生成できること。"""
        result = skill.execute(
            title="2カラムテスト",
            slides=json.dumps([TWO_COLUMN_SLIDE]),
            file_path=tmp_pptx,
        )
        assert "生成しました" in result

    def test_notes_added_to_slide(self, skill: PptxGenerationSkill, tmp_pptx: str) -> None:
        """発表者ノートがスライドに追加されること。"""
        slide_with_notes = {**TITLE_SLIDE, "notes": "これは発表者ノートです"}
        result = skill.execute(
            title="ノートテスト",
            slides=json.dumps([slide_with_notes]),
            file_path=tmp_pptx,
        )
        assert "生成しました" in result
        prs = Presentation(tmp_pptx)
        notes_tf = prs.slides[0].notes_slide.notes_text_frame
        assert "発表者ノート" in notes_tf.text

    def test_bullets_with_level(self, skill: PptxGenerationSkill, tmp_pptx: str) -> None:
        """ネストされた箇条書き（levelあり）が生成できること。"""
        nested_bullets = {
            "type": "bullets",
            "title": "ネスト箇条書き",
            "bullets": [
                {"level": 0, "text": "トップレベル"},
                {"level": 1, "text": "サブレベル1"},
                {"level": 2, "text": "サブレベル2"},
            ],
        }
        result = skill.execute(
            title="ネストテスト",
            slides=json.dumps([nested_bullets]),
            file_path=tmp_pptx,
        )
        assert "生成しました" in result

    def test_result_contains_slide_count(self, skill: PptxGenerationSkill, tmp_pptx: str) -> None:
        """結果メッセージにスライド数が含まれること。"""
        slides = [TITLE_SLIDE, BULLETS_SLIDE, CONTENT_SLIDE]
        result = skill.execute(
            title="スライド数テスト",
            slides=json.dumps(slides),
            file_path=tmp_pptx,
        )
        assert "3スライド" in result


# ── 異常系テスト ──────────────────────────────────────────────────────────────

class TestPptxGenerationSkillError:
    """異常系テスト。"""

    def test_invalid_json_slides(self, skill: PptxGenerationSkill, tmp_pptx: str) -> None:
        """不正なJSON文字列を渡した場合にエラーメッセージを返すこと。"""
        result = skill.execute(
            title="エラーテスト",
            slides="これはJSONではありません{{{",
            file_path=tmp_pptx,
        )
        assert "パースに失敗" in result
        assert not os.path.exists(tmp_pptx)

    def test_empty_slides_list(self, skill: PptxGenerationSkill, tmp_pptx: str) -> None:
        """空のスライドリストでもファイルが生成されること（0スライド）。"""
        result = skill.execute(
            title="空リストテスト",
            slides=json.dumps([]),
            file_path=tmp_pptx,
        )
        assert "生成しました" in result
        prs = Presentation(tmp_pptx)
        assert len(prs.slides) == 0

    def test_unknown_slide_type_falls_back_to_content(
        self, skill: PptxGenerationSkill, tmp_pptx: str
    ) -> None:
        """未知のスライドタイプはcontentとして処理されること。"""
        unknown_slide = {"type": "unknown_type", "title": "未知タイプ", "body": "本文"}
        result = skill.execute(
            title="未知タイプテスト",
            slides=json.dumps([unknown_slide]),
            file_path=tmp_pptx,
        )
        assert "生成しました" in result
        assert os.path.exists(tmp_pptx)

    def test_chart_slide_without_chart_def(
        self, skill: PptxGenerationSkill, tmp_pptx: str
    ) -> None:
        """chartフィールドが空でもエラーにならないこと。"""
        empty_chart = {"type": "chart", "title": "空チャート", "chart": {}}
        result = skill.execute(
            title="空チャートテスト",
            slides=json.dumps([empty_chart]),
            file_path=tmp_pptx,
        )
        assert "生成しました" in result

    def test_table_slide_without_headers(
        self, skill: PptxGenerationSkill, tmp_pptx: str
    ) -> None:
        """headersが空のテーブルスライドでもエラーにならないこと。"""
        empty_table = {"type": "table", "title": "空テーブル", "table": {"headers": [], "rows": []}}
        result = skill.execute(
            title="空テーブルテスト",
            slides=json.dumps([empty_table]),
            file_path=tmp_pptx,
        )
        assert "生成しました" in result

    def test_two_column_with_body_text(
        self, skill: PptxGenerationSkill, tmp_pptx: str
    ) -> None:
        """2カラムでbulletsの代わりにbodyを使っても生成できること。"""
        body_col = {
            "type": "two_column",
            "title": "bodyカラムテスト",
            "columns": [
                {"title": "左", "body": "左側の本文テキスト"},
                {"title": "右", "body": "右側の本文テキスト"},
            ],
        }
        result = skill.execute(
            title="bodyカラムテスト",
            slides=json.dumps([body_col]),
            file_path=tmp_pptx,
        )
        assert "生成しました" in result
