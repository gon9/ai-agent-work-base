"""
PptxJsGenerationSkillのテスト。

正常系・異常系を網羅する。
"""

from __future__ import annotations

import os
import textwrap

import pytest
from pptx import Presentation

from ai_agent_work_base.skills.pptxjs_generation import PptxJsGenerationSkill


@pytest.fixture
def skill() -> PptxJsGenerationSkill:
    """テスト用スキルインスタンス。"""
    return PptxJsGenerationSkill()


@pytest.fixture
def tmp_pptx(tmp_path) -> str:
    """一時出力ファイルパス。"""
    return str(tmp_path / "test.pptx")


# ── 最小限の動作確認スクリプト ────────────────────────────────────────────────

def _minimal_script(output_path: str) -> str:
    """最小限のpptxgenjsスクリプトを返す。"""
    return textwrap.dedent(f"""
        const PptxGenJS = require("pptxgenjs");
        const pres = new PptxGenJS();
        let slide = pres.addSlide();
        slide.addText("テストタイトル", {{ x: 1, y: 1, w: 8, h: 1, fontSize: 24, bold: true }});
        pres.writeFile({{ fileName: "{output_path}" }});
    """).strip()


def _multi_slide_script(output_path: str) -> str:
    """複数スライドのスクリプトを返す。"""
    return textwrap.dedent(f"""
        const PptxGenJS = require("pptxgenjs");
        const pres = new PptxGenJS();

        // スライド1: タイトル
        let slide1 = pres.addSlide();
        slide1.addShape(pres.ShapeType.rect, {{ x: 0, y: 0, w: "100%", h: 1.2, fill: {{ color: "1F497D" }}, line: {{ color: "1F497D" }} }});
        slide1.addText("インターネットの未来", {{ x: 0.5, y: 0.2, w: 9, h: 0.8, fontSize: 28, color: "FFFFFF", bold: true }});

        // スライド2: 箇条書き
        let slide2 = pres.addSlide();
        slide2.addText([
          {{ text: "• Web3の台頭", options: {{ fontSize: 16, color: "363636", breakLine: true }} }},
          {{ text: "• AIの深い統合", options: {{ fontSize: 16, color: "363636" }} }}
        ], {{ x: 0.5, y: 1.5, w: 9, h: 3 }});

        // スライド3: チャート
        let slide3 = pres.addSlide();
        slide3.addChart(pres.ChartType.bar,
          [{{ name: "成長率", labels: ["2022","2023","2024"], values: [10, 25, 40] }}],
          {{ x: 0.5, y: 1.5, w: 9, h: 4.5 }}
        );

        pres.writeFile({{ fileName: "{output_path}" }});
    """).strip()


def _script_with_codeblock(output_path: str) -> str:
    """コードブロック付きスクリプト（LLMが返す形式）を返す。"""
    inner = _minimal_script(output_path)
    return f"```javascript\n{inner}\n```"


# ── 正常系テスト ──────────────────────────────────────────────────────────────

class TestPptxJsGenerationSkillNormal:
    """正常系テスト。"""

    def test_skill_name(self, skill: PptxJsGenerationSkill) -> None:
        """スキル名が正しいこと。"""
        assert skill.name == "generate_pptx_js"

    def test_skill_description_not_empty(self, skill: PptxJsGenerationSkill) -> None:
        """説明が空でないこと。"""
        assert len(skill.description) > 0

    def test_parameters_schema(self, skill: PptxJsGenerationSkill) -> None:
        """パラメータスキーマにrequiredフィールドが含まれること。"""
        params = skill.parameters
        assert params["type"] == "object"
        assert "script" in params["properties"]
        assert "file_path" in params["properties"]
        assert set(params["required"]) == {"script", "file_path"}

    def test_generate_minimal_pptx(self, skill: PptxJsGenerationSkill, tmp_pptx: str) -> None:
        """最小限のスクリプトでPPTXが生成されること。"""
        script = _minimal_script(tmp_pptx)
        result = skill.execute(script=script, file_path=tmp_pptx)
        assert "生成しました" in result
        assert os.path.exists(tmp_pptx)

    def test_generated_file_is_valid_pptx(self, skill: PptxJsGenerationSkill, tmp_pptx: str) -> None:
        """生成されたファイルがpptxとして読み込めること。"""
        script = _minimal_script(tmp_pptx)
        skill.execute(script=script, file_path=tmp_pptx)
        prs = Presentation(tmp_pptx)
        assert len(prs.slides) >= 1

    def test_generate_multi_slide_pptx(self, skill: PptxJsGenerationSkill, tmp_pptx: str) -> None:
        """複数スライド（チャート含む）のPPTXが生成されること。"""
        script = _multi_slide_script(tmp_pptx)
        result = skill.execute(script=script, file_path=tmp_pptx)
        assert "生成しました" in result
        prs = Presentation(tmp_pptx)
        assert len(prs.slides) == 3

    def test_strips_markdown_codeblock(self, skill: PptxJsGenerationSkill, tmp_pptx: str) -> None:
        """LLMがコードブロックで返したスクリプトでも生成できること。"""
        script = _script_with_codeblock(tmp_pptx)
        result = skill.execute(script=script, file_path=tmp_pptx)
        assert "生成しました" in result
        assert os.path.exists(tmp_pptx)

    def test_output_directory_created(self, skill: PptxJsGenerationSkill, tmp_path) -> None:
        """出力ディレクトリが存在しなくても自動作成されること。"""
        nested_path = str(tmp_path / "nested" / "dir" / "output.pptx")
        script = _minimal_script(nested_path)
        result = skill.execute(script=script, file_path=nested_path)
        assert "生成しました" in result
        assert os.path.exists(nested_path)

    def test_result_contains_file_path(self, skill: PptxJsGenerationSkill, tmp_pptx: str) -> None:
        """結果メッセージにファイルパスが含まれること。"""
        script = _minimal_script(tmp_pptx)
        result = skill.execute(script=script, file_path=tmp_pptx)
        assert tmp_pptx in result or "test.pptx" in result

    def test_result_contains_file_size(self, skill: PptxJsGenerationSkill, tmp_pptx: str) -> None:
        """結果メッセージにファイルサイズが含まれること。"""
        script = _minimal_script(tmp_pptx)
        result = skill.execute(script=script, file_path=tmp_pptx)
        assert "KB" in result

    def test_inject_file_path_writeFile_object(self, skill: PptxJsGenerationSkill) -> None:
        """writeFile({fileName: ...})形式のパスが書き換えられること。"""
        script = 'pres.writeFile({ fileName: "old_path.pptx" });'
        result = skill._inject_file_path(script, "/new/path/output.pptx")
        assert "/new/path/output.pptx" in result
        assert "old_path.pptx" not in result

    def test_inject_file_path_writeFile_string(self, skill: PptxJsGenerationSkill) -> None:
        """writeFile("...")形式のパスが書き換えられること。"""
        script = 'pres.writeFile("old_path.pptx");'
        result = skill._inject_file_path(script, "/new/path/output.pptx")
        assert "/new/path/output.pptx" in result
        assert "old_path.pptx" not in result


# ── 異常系テスト ──────────────────────────────────────────────────────────────

class TestPptxJsGenerationSkillError:
    """異常系テスト。"""

    def test_invalid_js_syntax(self, skill: PptxJsGenerationSkill, tmp_pptx: str) -> None:
        """構文エラーのあるスクリプトはエラーメッセージを返すこと。"""
        broken_script = "const pres = new PptxGenJS(; // 構文エラー"
        result = skill.execute(script=broken_script, file_path=tmp_pptx)
        assert "エラー" in result
        assert not os.path.exists(tmp_pptx)

    def test_script_without_writefile(self, skill: PptxJsGenerationSkill, tmp_pptx: str) -> None:
        """writeFileを呼ばないスクリプトはファイル未生成メッセージを返すこと。"""
        no_write_script = textwrap.dedent("""
            const PptxGenJS = require("pptxgenjs");
            const pres = new PptxGenJS();
            let slide = pres.addSlide();
            slide.addText("test", { x: 1, y: 1, w: 8, h: 1 });
            // writeFileを呼ばない
        """).strip()
        result = skill.execute(script=no_write_script, file_path=tmp_pptx)
        assert "生成されませんでした" in result or "エラー" in result

    def test_runtime_error_in_script(self, skill: PptxJsGenerationSkill, tmp_pptx: str) -> None:
        """実行時エラーのあるスクリプトはエラーメッセージを返すこと。"""
        runtime_error_script = textwrap.dedent(f"""
            const PptxGenJS = require("pptxgenjs");
            const pres = new PptxGenJS();
            undefinedFunction();  // 実行時エラー
            pres.writeFile({{ fileName: "{tmp_pptx}" }});
        """).strip()
        result = skill.execute(script=runtime_error_script, file_path=tmp_pptx)
        assert "エラー" in result
