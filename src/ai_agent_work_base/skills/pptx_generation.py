"""
python-pptxを使用してリッチなPowerPointファイルを生成するスキル。

対応機能:
- タイトルスライド・コンテンツスライド・箇条書きスライド
- 棒グラフ・折れ線グラフ・円グラフの埋め込み
- テーブル（表）の埋め込み
- カスタムテーマカラー・フォント設定
- スライドノート（発表者ノート）
- 生成後にPowerPointで直接編集可能な.pptxファイル出力
"""

from __future__ import annotations

import json
import os
import re
from typing import Any

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Cm, Pt

from .base import BaseSkill

# ── デフォルトテーマカラー ──────────────────────────────────────────────────────
_THEME = {
    "primary": RGBColor(0x1F, 0x49, 0x7D),    # 濃紺
    "secondary": RGBColor(0x2E, 0x75, 0xB6),  # 青
    "accent": RGBColor(0xED, 0x7D, 0x31),     # オレンジ
    "light": RGBColor(0xD6, 0xE4, 0xF0),      # 薄青
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "dark": RGBColor(0x26, 0x26, 0x26),
    "gray": RGBColor(0x76, 0x76, 0x76),
}

# スライドサイズ (16:9 ワイド)
_SLIDE_W = Cm(33.87)
_SLIDE_H = Cm(19.05)

# チャートタイプのマッピング
_CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar_horizontal": XL_CHART_TYPE.BAR_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "line_markers": XL_CHART_TYPE.LINE_MARKERS,
    "pie": XL_CHART_TYPE.PIE,
    "doughnut": XL_CHART_TYPE.DOUGHNUT,
    "area": XL_CHART_TYPE.AREA,
    "scatter": XL_CHART_TYPE.XY_SCATTER,
}


def _set_run_font(run, size_pt: float, bold: bool = False, color: RGBColor | None = None) -> None:
    """ランのフォントを設定する。"""
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def _fill_shape(shape, color: RGBColor) -> None:
    """シェイプの背景色を設定する。"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def _add_header_bar(slide, title_text: str, subtitle_text: str = "") -> None:
    """スライド上部にカラーヘッダーバーとタイトルを追加する。"""
    # ヘッダー背景
    bar = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        0, 0, _SLIDE_W, Cm(3.2),
    )
    _fill_shape(bar, _THEME["primary"])
    bar.line.fill.background()

    # タイトルテキスト
    tf = bar.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title_text
    _set_run_font(run, 24, bold=True, color=_THEME["white"])

    # サブタイトル（あれば）
    if subtitle_text:
        sub_box = slide.shapes.add_textbox(Cm(1.0), Cm(2.6), Cm(28), Cm(0.9))
        tf2 = sub_box.text_frame
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle_text
        _set_run_font(run2, 12, color=_THEME["light"])


def _add_footer(slide, page_num: int, total: int, footer_text: str = "") -> None:
    """スライド下部にフッターとページ番号を追加する。"""
    # フッターライン
    line = slide.shapes.add_shape(1, 0, Cm(18.3), _SLIDE_W, Cm(0.05))
    _fill_shape(line, _THEME["secondary"])
    line.line.fill.background()

    # フッターテキスト
    if footer_text:
        ft_box = slide.shapes.add_textbox(Cm(1.0), Cm(18.4), Cm(25), Cm(0.6))
        tf = ft_box.text_frame
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = footer_text
        _set_run_font(run, 9, color=_THEME["gray"])

    # ページ番号
    pn_box = slide.shapes.add_textbox(Cm(30.5), Cm(18.4), Cm(3.0), Cm(0.6))
    tf = pn_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{page_num} / {total}"
    _set_run_font(run, 9, color=_THEME["gray"])


class PptxGenerationSkill(BaseSkill):
    """
    python-pptxを使用してリッチなPowerPointファイルを生成するスキル。

    スライド定義はJSON形式で受け取り、各スライドのタイプに応じて
    タイトル・コンテンツ・チャート・テーブルなどを生成する。
    生成された.pptxファイルはPowerPoint/LibreOfficeで直接編集可能。
    """

    @property
    def name(self) -> str:
        return "generate_pptx"

    @property
    def description(self) -> str:
        return (
            "JSON形式のスライド定義に基づいて、編集可能なPowerPoint(.pptx)ファイルを生成します。"
            "棒グラフ・折れ線グラフ・円グラフ・テーブルなどの図表も埋め込み可能です。"
        )

    @property
    def parameters(self) -> dict[str, Any]:
        return {
            "type": "object",
            "properties": {
                "title": {
                    "type": "string",
                    "description": "プレゼンテーションのタイトル",
                },
                "slides": {
                    "type": "string",
                    "description": (
                        "スライド定義のJSON文字列。各要素は以下のフィールドを持つ:\n"
                        "  type: 'title' | 'content' | 'bullets' | 'chart' | 'table' | 'two_column'\n"
                        "  title: スライドタイトル\n"
                        "  subtitle: サブタイトル (typeがtitleの場合)\n"
                        "  body: 本文テキスト (typeがcontentの場合)\n"
                        "  bullets: 箇条書きリスト (typeがbulletsの場合)\n"
                        "  notes: 発表者ノート\n"
                        "  chart: チャート定義 (typeがchartの場合)\n"
                        "    chart_type: 'bar'|'line'|'pie'|'doughnut'|'area'|'bar_horizontal'\n"
                        "    categories: カテゴリラベルのリスト\n"
                        "    series: [{name, values}] のリスト\n"
                        "  table: テーブル定義 (typeがtableの場合)\n"
                        "    headers: ヘッダー行のリスト\n"
                        "    rows: データ行のリスト（各行はリスト）\n"
                        "  columns: two_columnの場合の左右カラム定義 [{title, bullets}]\n"
                    ),
                },
                "file_path": {
                    "type": "string",
                    "description": "出力ファイルパス (例: output/presentation.pptx)",
                },
                "footer": {
                    "type": "string",
                    "description": "フッターに表示するテキスト（省略可）",
                },
            },
            "required": ["title", "slides", "file_path"],
        }

    def execute(
        self,
        title: str,
        slides: str | list,
        file_path: str,
        footer: str = "",
        **kwargs,
    ) -> str:
        """
        スライド定義に基づいてPowerPointファイルを生成する。

        Args:
            title: プレゼンテーションのタイトル
            slides: スライド定義のJSON文字列またはリスト
            file_path: 出力ファイルパス
            footer: フッターテキスト（省略可）

        Returns:
            生成結果メッセージ
        """
        try:
            slide_defs = self._parse_slides(slides)
        except (json.JSONDecodeError, ValueError) as e:
            return f"スライド定義のパースに失敗しました: {e}"

        try:
            prs = Presentation()
            prs.slide_width = _SLIDE_W
            prs.slide_height = _SLIDE_H

            # 空白レイアウト (index=6) を使用して完全カスタム描画
            blank_layout = prs.slide_layouts[6]

            total = len(slide_defs)
            for page_num, slide_def in enumerate(slide_defs, start=1):
                slide = prs.slides.add_slide(blank_layout)
                slide_type = slide_def.get("type", "content")

                if slide_type == "title":
                    self._build_title_slide(slide, slide_def, page_num, total, footer)
                elif slide_type == "bullets":
                    self._build_bullets_slide(slide, slide_def, page_num, total, footer)
                elif slide_type == "chart":
                    self._build_chart_slide(slide, slide_def, page_num, total, footer)
                elif slide_type == "table":
                    self._build_table_slide(slide, slide_def, page_num, total, footer)
                elif slide_type == "two_column":
                    self._build_two_column_slide(slide, slide_def, page_num, total, footer)
                else:
                    self._build_content_slide(slide, slide_def, page_num, total, footer)

                # 発表者ノート
                notes_text = slide_def.get("notes", "")
                if notes_text:
                    notes_slide = slide.notes_slide
                    tf = notes_slide.notes_text_frame
                    tf.text = notes_text

            os.makedirs(os.path.dirname(os.path.abspath(file_path)), exist_ok=True)
            prs.save(file_path)
            return f"PowerPointファイルを生成しました: {file_path} ({total}スライド)"

        except Exception as e:
            return f"PowerPointファイルの生成中にエラーが発生しました: {e}"

    # ── パース ────────────────────────────────────────────────────────────────

    def _parse_slides(self, slides: str | list) -> list[dict]:
        """スライド定義をパースしてリストで返す。"""
        if isinstance(slides, list):
            return slides
        # コードブロック除去
        cleaned = re.sub(r"^```[a-zA-Z]*\n?", "", slides.strip())
        cleaned = re.sub(r"\n?```$", "", cleaned.strip())
        return json.loads(cleaned)

    # ── スライドビルダー ──────────────────────────────────────────────────────

    def _build_title_slide(
        self, slide, slide_def: dict, page_num: int, total: int, footer: str
    ) -> None:
        """タイトルスライドを構築する。"""
        # 全面背景
        bg = slide.shapes.add_shape(1, 0, 0, _SLIDE_W, _SLIDE_H)
        _fill_shape(bg, _THEME["primary"])
        bg.line.fill.background()

        # アクセントライン
        accent = slide.shapes.add_shape(1, 0, Cm(11.5), _SLIDE_W, Cm(0.3))
        _fill_shape(accent, _THEME["accent"])
        accent.line.fill.background()

        # タイトル
        title_box = slide.shapes.add_textbox(Cm(2.5), Cm(5.5), Cm(28.5), Cm(5.0))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = slide_def.get("title", "")
        _set_run_font(run, 40, bold=True, color=_THEME["white"])

        # サブタイトル
        subtitle = slide_def.get("subtitle", "")
        if subtitle:
            sub_box = slide.shapes.add_textbox(Cm(2.5), Cm(12.2), Cm(28.5), Cm(2.5))
            tf2 = sub_box.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            run2 = p2.add_run()
            run2.text = subtitle
            _set_run_font(run2, 18, color=_THEME["light"])

        _add_footer(slide, page_num, total, footer)

    def _build_content_slide(
        self, slide, slide_def: dict, page_num: int, total: int, footer: str
    ) -> None:
        """コンテンツスライドを構築する。"""
        _add_header_bar(slide, slide_def.get("title", ""), slide_def.get("subtitle", ""))

        body = slide_def.get("body", "")
        if body:
            body_box = slide.shapes.add_textbox(Cm(1.5), Cm(3.8), Cm(30.5), Cm(13.5))
            tf = body_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = body
            _set_run_font(run, 16, color=_THEME["dark"])

        _add_footer(slide, page_num, total, footer)

    def _build_bullets_slide(
        self, slide, slide_def: dict, page_num: int, total: int, footer: str
    ) -> None:
        """箇条書きスライドを構築する。"""
        _add_header_bar(slide, slide_def.get("title", ""), slide_def.get("subtitle", ""))

        bullets: list = slide_def.get("bullets", [])
        body_box = slide.shapes.add_textbox(Cm(1.5), Cm(3.8), Cm(30.5), Cm(13.5))
        tf = body_box.text_frame
        tf.word_wrap = True

        for i, bullet in enumerate(bullets):
            if isinstance(bullet, dict):
                level = bullet.get("level", 0)
                text = bullet.get("text", "")
            else:
                level = 0
                text = str(bullet)

            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            p.level = level
            indent = "    " * level
            marker = "•" if level == 0 else "◦"
            run = p.add_run()
            run.text = f"{indent}{marker}  {text}"
            font_size = 18 - level * 2
            color = _THEME["dark"] if level == 0 else _THEME["gray"]
            _set_run_font(run, font_size, bold=(level == 0), color=color)

        _add_footer(slide, page_num, total, footer)

    def _build_chart_slide(
        self, slide, slide_def: dict, page_num: int, total: int, footer: str
    ) -> None:
        """チャートスライドを構築する。"""
        _add_header_bar(slide, slide_def.get("title", ""), slide_def.get("subtitle", ""))

        chart_def = slide_def.get("chart", {})
        chart_type_key = chart_def.get("chart_type", "bar")
        chart_type = _CHART_TYPE_MAP.get(chart_type_key, XL_CHART_TYPE.COLUMN_CLUSTERED)

        categories = chart_def.get("categories", [])
        series_list = chart_def.get("series", [])

        chart_data = CategoryChartData()
        chart_data.categories = categories
        for s in series_list:
            chart_data.add_series(s.get("name", ""), s.get("values", []))

        chart_frame = slide.shapes.add_chart(
            chart_type,
            Cm(1.5),
            Cm(3.8),
            Cm(30.5),
            Cm(13.5),
            chart_data,
        )
        chart = chart_frame.chart

        # チャートスタイル調整
        chart.has_legend = len(series_list) > 1
        if chart.has_legend:
            chart.legend.position = 2  # BOTTOM
            chart.legend.include_in_layout = False

        _add_footer(slide, page_num, total, footer)

    def _build_table_slide(
        self, slide, slide_def: dict, page_num: int, total: int, footer: str
    ) -> None:
        """テーブルスライドを構築する。"""
        _add_header_bar(slide, slide_def.get("title", ""), slide_def.get("subtitle", ""))

        table_def = slide_def.get("table", {})
        headers: list[str] = table_def.get("headers", [])
        rows: list[list] = table_def.get("rows", [])

        if not headers:
            _add_footer(slide, page_num, total, footer)
            return

        num_cols = len(headers)
        num_rows = len(rows) + 1  # ヘッダー行含む

        col_width = Cm(30.5) // num_cols
        row_height = min(Cm(1.5), Cm(13.5) // num_rows)

        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            Cm(1.5), Cm(3.8),
            Cm(30.5), row_height * num_rows,
        )
        tbl = table_shape.table

        # 列幅均等
        for col_idx in range(num_cols):
            tbl.columns[col_idx].width = col_width

        # ヘッダー行
        for col_idx, header in enumerate(headers):
            cell = tbl.cell(0, col_idx)
            cell.text = str(header)
            cell.fill.solid()
            cell.fill.fore_color.rgb = _THEME["primary"]
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0] if p.runs else p.add_run()
            _set_run_font(run, 13, bold=True, color=_THEME["white"])

        # データ行
        for row_idx, row_data in enumerate(rows):
            bg_color = _THEME["light"] if row_idx % 2 == 0 else _THEME["white"]
            for col_idx, cell_val in enumerate(row_data):
                cell = tbl.cell(row_idx + 1, col_idx)
                cell.text = str(cell_val)
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg_color
                p = cell.text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.runs[0] if p.runs else p.add_run()
                _set_run_font(run, 12, color=_THEME["dark"])

        _add_footer(slide, page_num, total, footer)

    def _build_two_column_slide(
        self, slide, slide_def: dict, page_num: int, total: int, footer: str
    ) -> None:
        """2カラムスライドを構築する。"""
        _add_header_bar(slide, slide_def.get("title", ""), slide_def.get("subtitle", ""))

        columns: list[dict] = slide_def.get("columns", [{}, {}])
        col_w = Cm(14.5)
        col_gap = Cm(1.0)
        col_y = Cm(3.8)
        col_h = Cm(13.5)

        for i, col_def in enumerate(columns[:2]):
            col_x = Cm(1.5) + i * (col_w + col_gap)

            # カラムヘッダー
            col_title = col_def.get("title", "")
            if col_title:
                hdr = slide.shapes.add_shape(1, col_x, col_y, col_w, Cm(0.9))
                _fill_shape(hdr, _THEME["secondary"])
                hdr.line.fill.background()
                tf = hdr.text_frame
                p = tf.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = col_title
                _set_run_font(run, 14, bold=True, color=_THEME["white"])

            # カラムコンテンツ
            content_y = col_y + Cm(1.0)
            content_h = col_h - Cm(1.0)
            body_box = slide.shapes.add_textbox(col_x, content_y, col_w, content_h)
            tf2 = body_box.text_frame
            tf2.word_wrap = True

            bullets = col_def.get("bullets", [])
            body_text = col_def.get("body", "")

            if bullets:
                for j, bullet in enumerate(bullets):
                    text = bullet if isinstance(bullet, str) else bullet.get("text", "")
                    p2 = tf2.paragraphs[0] if j == 0 else tf2.add_paragraph()
                    run2 = p2.add_run()
                    run2.text = f"•  {text}"
                    _set_run_font(run2, 14, color=_THEME["dark"])
            elif body_text:
                p2 = tf2.paragraphs[0]
                run2 = p2.add_run()
                run2.text = body_text
                _set_run_font(run2, 14, color=_THEME["dark"])

        _add_footer(slide, page_num, total, footer)
